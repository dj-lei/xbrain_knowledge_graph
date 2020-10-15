from docx import Document
from pptx import Presentation
from abc import abstractmethod
import pandas as pd
import numpy as np
import re
import zipfile
import win32com
import win32com.client
import os
import platform


def Doctable(ls, row, column):
    df = pd.DataFrame(np.array(ls).reshape(row, column))  # reshape to the table shape
    df.columns = df.loc[0, :].values
    df = df.loc[1:, :].dropna(how="all").drop_duplicates().reset_index(drop=True)
    return df


class BaseDocHandle(object):

    def __init__(self, path):
        self.path = path
        self.proper_nouns = set()
        self.document = Document(self.path)
        # with open('kg_extract/google-10000-english-usa-no-swears.txt', 'r') as f:
        #     self.common_words = [token.strip() for token in f]

    def get_proper_nouns_csv(self):
        self.proper_nouns = set(pd.read_csv('proper_nouns.csv')['proper_nouns_name'].values)

    def get_docx_structure(self):
        count = 0
        doc_name = self.document.core_properties.title
        # self.get_proper_nouns_csv()

        level = 0
        content = [0 for _ in range(0, 10)]
        content[0] = doc_name
        temp_name = ''
        data = pd.DataFrame([], columns=['s', 'p', 'o'])
        for i in range(0, len(self.document.paragraphs)):
            doc = self.document.paragraphs[i].text.strip()
            style_name = self.document.paragraphs[i].style.name
            if doc in ['', '\n']:
                continue

            # key_words = str(list(self.proper_nouns.intersection(set(doc.split(' ')))))
            if style_name == 'Heading':
                data.loc[count, 's'] = content[0]
                data.loc[count, 'p'] = 'Heading 0'
                data.loc[count, 'o'] = doc
                # data.loc[count, 'key_words'] = key_words
                temp_name = doc
            elif 'Heading' in style_name:
                level = int(style_name.split(' ')[1])
                content[level] = doc

                data.loc[count, 's'] = content[level-1]
                data.loc[count, 'p'] = style_name
                data.loc[count, 'o'] = doc
                # data.loc[count, 'key_words'] = key_words
            else:
                if level == 0:
                    data.loc[count, 's'] = temp_name
                    data.loc[count, 'p'] = 'Heading ' + str(level + 1)
                    data.loc[count, 'o'] = doc
                    # data.loc[count, 'key_words'] = key_words
                    count = count + 1
                    continue

                if style_name in ['Caption', 'CaptionFigure']:
                    data.loc[count, 's'] = content[level]
                    data.loc[count, 'p'] = 'Caption'
                    data.loc[count, 'o'] = doc
                    # data.loc[count, 'key_words'] = key_words
                    count = count + 1
                    continue

                data.loc[count, 's'] = content[level]
                data.loc[count, 'p'] = 'Heading '+str(level+1)
                data.loc[count, 'o'] = doc
                # data.loc[count, 'key_words'] = key_words
            count = count + 1

        data = data.reset_index(drop=True)
        return data

    def get_docx_images(self, dstpath):
        doc = zipfile.ZipFile(self.path)
        for info in doc.infolist():
            if info.filename.endswith((".png", ".jpeg", ".jpg", ".emf")):
                if int(info.file_size) > 2000:
                    doc.extract(info.filename, dstpath)
        doc.close()

    def get_docx_tables(self):
        result = []
        for table in self.document.tables:
            ls = []
            if len(table.rows) <= 1:
                continue
            for row in table.rows:
                for cell in row.cells:
                    temp = []
                    for paragraph in cell.paragraphs:
                        temp.append(paragraph.text)
                    ls.append('\n'.join(temp))

            result.append(Doctable(ls, len(table.rows), len(table.rows[0].cells)))
        return result

    @abstractmethod
    def cut_special_symbols_and_filter_rules(self, word):
        pass

    @abstractmethod
    def get_proper_nouns(self):
        pass


class BasePptxHandle(object):
    def __init__(self, path):
        if platform.platform().startswith("Windows"):
            self.pwd = os.getcwd() + "\\"
        else:
            self.pwd = ''
        self.path = path

    def get_pptx_structure(self):
        prs = Presentation(self.pwd + self.path)

        slide = dict()
        for shape_num, shape in enumerate(prs.slides[9].shapes):
            data = dict()
            if shape.has_text_frame:
                data['type'] = 'text'
                content = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        content.append([run.text, run.font.size])
                data['content'] = content
            elif 'picture' in str(shape):
                data['type'] = 'image'
                data['content'] = shape.image.blob
            elif shape.has_table:
                data['type'] = 'table'
                ls = []
                for row in shape.table.rows:
                    for cell in row.cells:
                        temp = []
                        for paragraph in cell.text_frame.paragraphs:
                            temp.append(paragraph.text)
                        ls.append('\n'.join(temp))
                data['content'] = Doctable(ls, len(shape.table.rows), len(shape.table.rows[0].cells))
            else:
                continue
            data['shape_num'] = shape_num
            slide['shape_num_' + str(shape_num)] = data
        return slide

    def get_pptx_images(self, dstpath):
        application = win32com.client.Dispatch("PowerPoint.Application")
        presentation = application.Presentations.Open(self.pwd + self.path, WithWindow=False)
        for slide in presentation.Slides:
            slide.Export(self.pwd + dstpath, "JPG")
        application.Quit()


class RuDocHandle(BaseDocHandle):

    def __init__(self, path):
        super(RuDocHandle, self).__init__(path)

    def cut_special_symbols_and_filter_rules(self, word):
        r_symbols = "[\s\.\!\$%^*\"\'\:\;\,\?\(\)\[\]\“\”\‘\’\、]"
        r_symbols_special = '[\-\/]$'
        r_filter = '[0-9A-Z]-[0-9A-Z]'
        # cut_special_symbols
        self.word = re.sub(r_symbols, '', word).strip()
        self.word = re.sub(r_symbols_special, '', self.word)
        # filter_rules
        if len(re.findall(r_filter, self.word)) > 0:
            return ''
        return self.word

    def get_proper_nouns(self):
        count = 0
        proper_nouns_name = ''
        data = pd.DataFrame([], columns=['proper_nouns_name', 'text'])
        for i in range(0, len(self.document.paragraphs)):
            doc = self.document.paragraphs[i].text.strip()
            if doc not in ['', '\n']:
                for word in doc.split(' '):
                    if word.isupper():
                        proper_nouns_name = word
                    elif len(re.findall('[\(](.*?)[\)]', word)) > 0:
                        proper_nouns_name = re.findall('[\(](.*?)[\)]', word)[0]
                    else:
                        continue

                    proper_nouns_name = self.cut_special_symbols_and_filter_rules(proper_nouns_name)
                    if len(proper_nouns_name) <= 1:
                        continue
                    data.loc[count, 'proper_nouns_name'] = proper_nouns_name
                    data.loc[count, 'text'] = doc
                    count = count + 1
        data['file_name'] = self.document.core_properties.title
        data = data.sort_values(by='proper_nouns_name').drop_duplicates(['proper_nouns_name']).reset_index(drop=True)

        return data
