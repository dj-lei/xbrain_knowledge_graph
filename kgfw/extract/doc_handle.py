import re
import os
import win32com
import platform
from docx.document import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
from pptx import Presentation
from abc import abstractmethod
import pandas as pd
from io import BytesIO
from PIL import Image
from docx.shared import Pt
import numpy as np
import docx


def get_image_with_rel(doc, rid):
    for rel in doc.part._rels:
        rel = doc.part._rels[rel]
        if rel.rId == rid:
            return rel.target_part.blob


def iter_block_items(parent):
    if isinstance(parent, Document):
        parent_elm = parent.element.body
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            yield Table(child, parent)


def doctable(ls, row, column):
    df = pd.DataFrame(np.array(ls).reshape(row, column))  # reshape to the table shape
    df.columns = df.loc[0, :].values
    df = df.loc[1:, :].dropna(how="all").drop_duplicates().reset_index(drop=True)
    return df


def genarate_table(table):
    ls = []
    for row in table.rows:
        for cell in row.cells:
            temp = []
            for paragraph in cell.paragraphs:
                temp.append(paragraph.text)
            ls.append('\n'.join(temp))

    return doctable(ls, len(table.rows), len(table.rows[0].cells))


class BaseDocxHandle(object):

    def __init__(self, path):
        self.path = path
        self.proper_nouns = set()
        self.document = docx.Document(self.path)
        # with open('extract/google-10000-english-usa-no-swears.txt', 'r') as f:
        #     self.common_words = [token.strip() for token in f]

    def get_docx_structure(self):
        tmp = []
        for para in iter_block_items(self.document):
            style_name = para.style.name

            if 'docx.table.Table' in str(para):
                table = genarate_table(para)
                if len(table) < 1:
                    continue
                tmp.append({'type': 'table', 'style': style_name, 'content': table})
                continue

            if 'imagedata' in para._p.xml:
                rid = re.findall('imagedata r:id=(.*?) ', para._p.xml)[0].replace('"', '')
                tmp.append({'type': 'image', 'style': style_name, 'content': get_image_with_rel(self.document, rid)})
                continue

            doc = para.text.strip()
            if doc in ['', '\n']:
                continue
            if style_name in ['List abc double line', 'List number single line']:
                style_name = 'List Bullet'
            tmp.append({'type': 'text', 'style': style_name, 'content': doc})
        return {'name': self.document.core_properties.title, 'content': tmp}

    def get_catalog(self):
        data = self.get_docx_structure()
        prev_level = 0
        content = [0 for _ in range(0, 10)]
        count = [0 for _ in range(0, 10)]
        tmp = []
        for elm in data['content']:
            if 'Heading ' in elm['style'].strip():
                current_level = int(elm['style'].strip().split(' ')[1])

                if prev_level == current_level:
                    content[current_level - 1] = elm['content']
                    count[current_level - 1] += 1
                elif prev_level + 1 == current_level:
                    content[current_level - 1] = elm['content']
                    prev_level = current_level
                    count[current_level - 1] = 1
                elif prev_level > current_level:
                    content[current_level - 1] = elm['content']
                    prev_level = current_level
                    count[current_level - 1] += 1
                    content[current_level:] = [0 for _ in range(current_level, 10)]
                    count[current_level:] = [0 for _ in range(current_level, 10)]
                tmp.append({'chapter': '.'.join([str(v) for v in count if v != 0]), 'title': [str(v) for v in content if v != 0][-1], 'path_content':
                            ' '.join([str(v) for v in content if v != 0])})
        return {'name': self.document.core_properties.title, 'catalog': tmp}

    def recovery_docx(self):
        data = self.get_docx_structure()
        catalog = self.get_catalog()

        document = docx.Document()
        document.add_heading(data['name'], 0)
        count = 0

        for elm in data['content']:
            if elm['type'] == 'text':
                if 'Heading ' in elm['style']:
                    run = document.add_heading().add_run(catalog['catalog'][count]['chapter'] + ' ' + elm['content'])
                    font = run.font
                    font.size = Pt(22 - int(elm['style'].split(' ')[1]) * 2)
                    count += 1
                else:
                    if 'List' in elm['style']:
                        run = document.add_paragraph(style=elm['style']).add_run(elm['content'])
                    else:
                        run = document.add_paragraph().add_run(elm['content'])
                    font = run.font
                    font.size = Pt(12)
                font.name = 'Arial'

            if elm['type'] == 'image':
                pImage = Image.open(BytesIO(bytes(elm['content'])))
                pImage = pImage.resize((450, 300), Image.ANTIALIAS)
                pImage.save('temp.png')
                document.add_picture('temp.png')

            if elm['type'] == 'table':
                tal = elm['content']
                table = document.add_table(rows=1, cols=len(tal.columns))
                table.style = 'TableGrid'
                hdr_cells = table.rows[0].cells
                for i, column in enumerate(tal.columns):
                    hdr_cells[i].text = column
                for values in tal.values:
                    row_cells = table.add_row().cells
                    for j, value in enumerate(values):
                        row_cells[j].text = str(value)

        document.save('demo.docx')

    # def get_docx_structure_v2(self):
    #     count = 0
    #     doc_name = self.document.core_properties.title
    #     # self.get_proper_nouns_csv()
    #
    #     level = 0
    #     content = [0 for _ in range(0, 10)]
    #     content[0] = doc_name
    #     temp_name = ''
    #     data = pd.DataFrame([], columns=['s', 'p', 'o'])
    #     for i in range(0, len(self.document.paragraphs)):
    #         doc = self.document.paragraphs[i].text.strip()
    #         style_name = self.document.paragraphs[i].style.name
    #         if doc in ['', '\n']:
    #             continue
    #
    #         # key_words = str(list(self.proper_nouns.intersection(set(doc.split(' ')))))
    #         if style_name == 'Heading':
    #             data.loc[count, 's'] = content[0]
    #             data.loc[count, 'p'] = 'Heading 0'
    #             data.loc[count, 'o'] = doc
    #             # data.loc[count, 'key_words'] = key_words
    #             temp_name = doc
    #         elif 'Heading' in style_name:
    #             level = int(style_name.split(' ')[1])
    #             content[level] = doc
    #
    #             data.loc[count, 's'] = content[level-1]
    #             data.loc[count, 'p'] = style_name
    #             data.loc[count, 'o'] = doc
    #             # data.loc[count, 'key_words'] = key_words
    #         else:
    #             if level == 0:
    #                 data.loc[count, 's'] = temp_name
    #                 data.loc[count, 'p'] = 'Heading ' + str(level + 1)
    #                 data.loc[count, 'o'] = doc
    #                 # data.loc[count, 'key_words'] = key_words
    #                 count = count + 1
    #                 continue
    #
    #             if style_name in ['Caption', 'CaptionFigure']:
    #                 data.loc[count, 's'] = content[level]
    #                 data.loc[count, 'p'] = 'Caption'
    #                 data.loc[count, 'o'] = doc
    #                 # data.loc[count, 'key_words'] = key_words
    #                 count = count + 1
    #                 continue
    #
    #             data.loc[count, 's'] = content[level]
    #             data.loc[count, 'p'] = 'Heading '+str(level+1)
    #             data.loc[count, 'o'] = doc
    #             # data.loc[count, 'key_words'] = key_words
    #         count = count + 1
    #
    #     data = data.reset_index(drop=True)
    #     return data


class BasePptxHandle(object):
    def __init__(self, path):
        if platform.platform().startswith("Windows"):
            self.pwd = os.getcwd() + "\\"
        else:
            self.pwd = ''
        self.path = path

    def get_pptx_structure(self):
        prs = Presentation(self.pwd + self.path)
        tmp = []
        for slide_num, slide_s in enumerate(prs.slides):
            slide = []
            for shape_num, shape in enumerate(slide_s.shapes):
                data = dict()
                if shape.has_text_frame:
                    data['type'] = 'text'
                    content = []
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            content.append([run.text, run.font.size])
                    data['content'] = content
                elif 'picture' in str(shape):
                    data['type'] = 'image'
                    data['content'] = shape.image.blob
                elif shape.has_table:
                    data['type'] = 'table'
                    ls = []
                    for row in shape.table.rows:
                        for cell in row.cells:
                            temp = []
                            for paragraph in cell.text_frame.paragraphs:
                                temp.append(paragraph.text)
                            ls.append('\n'.join(temp))
                    data['content'] = doctable(ls, len(shape.table.rows), len(shape.table.rows[0].cells))
                else:
                    continue
                slide.append(data)
            tmp.append({'page': slide_num, 'slide': slide})
        return {'name': prs.core_properties.title, 'slides': tmp}

    def export_pptx_images(self):
        application = win32com.client.Dispatch("PowerPoint.Application")
        presentation = application.Presentations.Open(self.pwd + self.path, WithWindow=False)
        for slide in presentation.Slides:
            slide.Export(self.pwd + r"tmp/tmp", "JPG")
            # other handle
        application.Quit()


# class RuDocHandle(BaseDocxHandle):
#
#     def __init__(self, path):
#         super(RuDocHandle, self).__init__(path)
#
#     def cut_special_symbols_and_filter_rules(self, word):
#         r_symbols = "[\s\.\!\$%^*\"\'\:\;\,\?\(\)\[\]\“\”\‘\’\、]"
#         r_symbols_special = '[\-\/]$'
#         r_filter = '[0-9A-Z]-[0-9A-Z]'
#         # cut_special_symbols
#         self.word = re.sub(r_symbols, '', word).strip()
#         self.word = re.sub(r_symbols_special, '', self.word)
#         # filter_rules
#         if len(re.findall(r_filter, self.word)) > 0:
#             return ''
#         return self.word
#
#     def get_proper_nouns(self):
#         count = 0
#         proper_nouns_name = ''
#         data = pd.DataFrame([], columns=['proper_nouns_name', 'text'])
#         for i in range(0, len(self.document.paragraphs)):
#             doc = self.document.paragraphs[i].text.strip()
#             if doc not in ['', '\n']:
#                 for word in doc.split(' '):
#                     if word.isupper():
#                         proper_nouns_name = word
#                     elif len(re.findall('[\(](.*?)[\)]', word)) > 0:
#                         proper_nouns_name = re.findall('[\(](.*?)[\)]', word)[0]
#                     else:
#                         continue
#
#                     proper_nouns_name = self.cut_special_symbols_and_filter_rules(proper_nouns_name)
#                     if len(proper_nouns_name) <= 1:
#                         continue
#                     data.loc[count, 'proper_nouns_name'] = proper_nouns_name
#                     data.loc[count, 'text'] = doc
#                     count = count + 1
#         data['file_name'] = self.document.core_properties.title
#         data = data.sort_values(by='proper_nouns_name').drop_duplicates(['proper_nouns_name']).reset_index(drop=True)
#
#         return data
